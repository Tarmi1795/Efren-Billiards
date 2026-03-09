from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import copy

# ── Color Palette (Premium Dark Gold) ──────────────────────────────────────
DARK_BG     = RGBColor(0x0A, 0x0A, 0x0C)   # Almost black
DARK_CARD   = RGBColor(0x12, 0x12, 0x16)   # Dark card
GOLD        = RGBColor(0xC5, 0xA0, 0x59)   # Brand gold
GOLD_LIGHT  = RGBColor(0xE8, 0xCB, 0x8A)   # Light gold
MAROON      = RGBColor(0x7B, 0x1F, 0x2A)   # Deep maroon
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0x9C, 0x9C, 0xA8)
GRAY_LIGHT  = RGBColor(0xD0, 0xD0, 0xD8)

W = Inches(13.33)
H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def gold_bar(slide, top_offset=Inches(0.08)):
    """Thin gold accent bar at the top."""
    add_rect(slide, 0, top_offset, W, Inches(0.05), GOLD)


def slide_number(slide, num, total):
    add_text(slide, f"{num:02d} / {total:02d}",
             W - Inches(1.4), H - Inches(0.45), Inches(1.2), Inches(0.35),
             font_size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def section_label(slide, label):
    """Top-right section indicator."""
    add_text(slide, label.upper(),
             W - Inches(3.2), Inches(0.15), Inches(3.0), Inches(0.35),
             font_size=8, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)


def bullet_block(slide, items, left, top, width, spacing=Inches(0.42), font_size=14, color=GRAY_LIGHT):
    for item in items:
        icon = "▸ "
        add_text(slide, icon + item, left, top, width, Inches(0.38),
                 font_size=font_size, color=color)
        top += spacing
    return top


# ── Build Presentation ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # Completely blank
TOTAL = 12

slides_data = []
# We'll build each slide imperatively


# ── SLIDE 1: Cover ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)

# Left gradient bar
add_rect(s, 0, 0, Inches(0.18), H, MAROON)
add_rect(s, Inches(0.18), 0, Inches(0.06), H, GOLD)

# Gold horizontal line mid-page
add_rect(s, Inches(1.2), Inches(3.55), Inches(11), Inches(0.03), GOLD)

# Title
add_text(s, "EFREN BILLIARDS & EVENTS", Inches(1.2), Inches(1.0), Inches(11), Inches(1.2),
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "DOHA, QATAR  🎱", Inches(1.2), Inches(2.1), Inches(9), Inches(0.8),
         font_size=28, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

# Sub
add_text(s, "Platform Overview & Admin CMS Guide",
         Inches(1.2), Inches(3.7), Inches(9), Inches(0.55),
         font_size=18, color=GRAY_LIGHT, italic=True)

# Bottom line
add_text(s, "Powered by React · TypeScript · Supabase",
         Inches(1.2), H - Inches(0.8), Inches(9), Inches(0.4),
         font_size=11, color=GRAY, italic=True)
slide_number(s, 1, TOTAL)


# ── SLIDE 2: Project Overview ──────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
gold_bar(s)
section_label(s, "Overview")

add_text(s, "WHAT IS THIS PLATFORM?", Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         font_size=32, bold=True, color=WHITE)
add_rect(s, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.05), GOLD)

add_text(s, "A premium, full-stack web application built for Efren Billiards & Events — a world-class billiards club in Doha, Qatar. It manages tournaments, membership, digital content, and event booking from one unified platform.",
         Inches(0.8), Inches(1.4), Inches(11.6), Inches(1.2),
         font_size=16, color=GRAY_LIGHT, wrap=True)

# 3 stat boxes
for i, (label, val, sub) in enumerate([
    ("🏆", "Tournaments", "Brackets + Elimination"),
    ("🛠️", "Admin CMS", "12-Module Control Panel"),
    ("📊", "Live Rankings", "Billiards · Darts · Chess"),
]):
    bx = Inches(0.8) + i * Inches(4.1)
    add_rect(s, bx, Inches(2.9), Inches(3.7), Inches(2.0), DARK_CARD)
    add_rect(s, bx, Inches(2.9), Inches(3.7), Inches(0.06), GOLD)
    add_text(s, label, bx + Inches(0.2), Inches(3.05), Inches(0.6), Inches(0.6), font_size=28)
    add_text(s, val, bx + Inches(0.2), Inches(3.55), Inches(3.3), Inches(0.45),
             font_size=16, bold=True, color=WHITE)
    add_text(s, sub, bx + Inches(0.2), Inches(3.95), Inches(3.3), Inches(0.4),
             font_size=11, color=GOLD)

slide_number(s, 2, TOTAL)


# ── SLIDE 3: Tech Stack ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
gold_bar(s)
section_label(s, "Tech Stack")

add_text(s, "TECHNOLOGY STACK", Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         font_size=32, bold=True, color=WHITE)
add_rect(s, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.05), GOLD)

cols = [
    ("⚛️  Frontend", ["React 19", "TypeScript", "Tailwind CSS", "Framer Motion", "Lucide Icons"]),
    ("🗄️  Backend", ["Supabase (Postgres)", "Row-Level Security", "Auth (OTP + Google)", "Storage Buckets", "Real-time API"]),
    ("🛠️  Tooling", ["Vite", "npm", "ESLint", "Git / GitHub", "python-pptx (Docs)"]),
]
for i, (title, items) in enumerate(cols):
    bx = Inches(0.6) + i * Inches(4.2)
    add_rect(s, bx, Inches(1.55), Inches(3.9), Inches(5.0), DARK_CARD)
    add_text(s, title, bx + Inches(0.25), Inches(1.7), Inches(3.6), Inches(0.5),
             font_size=15, bold=True, color=GOLD)
    add_rect(s, bx + Inches(0.25), Inches(2.15), Inches(0.8), Inches(0.04), MAROON)
    top = Inches(2.35)
    for it in items:
        add_text(s, "•  " + it, bx + Inches(0.25), top, Inches(3.4), Inches(0.38),
                 font_size=13, color=GRAY_LIGHT)
        top += Inches(0.4)

slide_number(s, 3, TOTAL)


# ── SLIDE 4: Site Map ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
gold_bar(s)
section_label(s, "Pages")

add_text(s, "SITE MAP & PAGES", Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         font_size=32, bold=True, color=WHITE)
add_rect(s, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.05), GOLD)

pages = [
    ("🏠", "#home", "Homepage – Hero, Offerings, Events, Gallery, Timetable"),
    ("🎱", "#billiards", "Billiards Page – Legend Content + Live Rankings"),
    ("🎯", "#darts", "Darts Page – Precision Focus + Rankings"),
    ("♟️", "#chess", "Chess Page – Strategy Content + Rankings"),
    ("🏛️", "#event-place", "Premium Event Place – Budget Estimator + Gallery"),
    ("🎤", "#karaoke", "Karaoke – Private Room Showcase"),
    ("☕", "#coffee-menu", "Coffee Menu – Live Price Digital Menu"),
    ("🏆", "#tournaments", "Tournaments – Live Bracket Viewer + Registration"),
    ("👑", "#membership-packages", "Membership – Tiered Benefits Landing Page"),
    ("🛡️", "#admin-cms", "Admin CMS – 12-Module Content Management Portal"),
]

col1 = pages[:5]
col2 = pages[5:]

for i, (icon, route, desc) in enumerate(col1):
    top = Inches(1.55) + i * Inches(0.95)
    add_rect(s, Inches(0.5), top, Inches(6.1), Inches(0.82), DARK_CARD)
    add_text(s, icon + "  " + route, Inches(0.7), top + Inches(0.08), Inches(2.2), Inches(0.35),
             font_size=11, bold=True, color=GOLD)
    add_text(s, desc, Inches(0.7), top + Inches(0.4), Inches(5.6), Inches(0.35),
             font_size=11, color=GRAY_LIGHT)

for i, (icon, route, desc) in enumerate(col2):
    top = Inches(1.55) + i * Inches(0.95)
    add_rect(s, Inches(6.9), top, Inches(6.1), Inches(0.82), DARK_CARD)
    add_text(s, icon + "  " + route, Inches(7.1), top + Inches(0.08), Inches(2.5), Inches(0.35),
             font_size=11, bold=True, color=GOLD)
    add_text(s, desc, Inches(7.1), top + Inches(0.4), Inches(5.6), Inches(0.35),
             font_size=11, color=GRAY_LIGHT)

slide_number(s, 4, TOTAL)


# ── SLIDE 5: Admin CMS Overview ────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
gold_bar(s)
section_label(s, "CMS")

add_text(s, "ADMIN CMS PORTAL", Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         font_size=32, bold=True, color=WHITE)
add_rect(s, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.05), GOLD)
add_text(s, "A full 12-module management system. No coding required — change any content on the live site from a single panel.",
         Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.6), font_size=15, color=GRAY_LIGHT)

modules = [
    ("🎨", "Branding & Visuals", "Hero · Site Images · Gallery · Videos · Visual Tour"),
    ("🍽️", "Services & Menu", "Food Menu · Event Pricing · Match My Game · Membership Plans"),
    ("📅", "Operations & Contact", "Weekly Schedule · Contact Info · Social Links"),
    ("🏆", "Club Management", "Tournaments · Players · Rankings · Members · Settings"),
]

for i, (icon, title, sub) in enumerate(modules):
    row = i // 2
    col = i % 2
    bx = Inches(0.55) + col * Inches(6.4)
    by = Inches(2.2) + row * Inches(2.1)
    add_rect(s, bx, by, Inches(6.0), Inches(1.85), DARK_CARD)
    add_rect(s, bx, by, Inches(0.08), Inches(1.85), GOLD)
    add_text(s, icon + "  " + title, bx + Inches(0.25), by + Inches(0.2), Inches(5.5), Inches(0.5),
             font_size=17, bold=True, color=WHITE)
    add_text(s, sub, bx + Inches(0.25), by + Inches(0.72), Inches(5.5), Inches(0.8),
             font_size=12, color=GOLD)

slide_number(s, 5, TOTAL)


# ── SLIDE 6: Branding & Visuals Deep Dive ─────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
gold_bar(s)
section_label(s, "CMS — Branding")

add_text(s, "🎨  BRANDING & VISUALS", Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         font_size=30, bold=True, color=WHITE)
add_rect(s, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.05), GOLD)

modules_detail = [
    ("Hero Section", "Edit homepage headline, sub-heading, and CTA button.\nChanges are saved to the database and sync instantly across all devices."),
    ("Site Images", "Swap logo, hero backgrounds, and section photos by pasting direct image URLs."),
    ("Image Gallery", "Upload and caption photos. Displayed in the scrolling gallery on the homepage."),
    ("YouTube Videos", "Paste any YouTube link. The system auto-embeds it into the video player section."),
    ("Visual Tour", "Assign descriptions and images to each area of the club (e.g., Main Hall, VIP rooms)."),
]

top = Inches(1.5)
for title, desc in modules_detail:
    add_rect(s, Inches(0.6), top, Inches(12.1), Inches(0.85), DARK_CARD)
    add_rect(s, Inches(0.6), top, Inches(0.07), Inches(0.85), GOLD)
    add_text(s, title, Inches(0.85), top + Inches(0.05), Inches(3.5), Inches(0.38),
             font_size=13, bold=True, color=GOLD_LIGHT)
    add_text(s, desc, Inches(4.0), top + Inches(0.05), Inches(8.5), Inches(0.65),
             font_size=11, color=GRAY_LIGHT)
    top += Inches(0.97)

slide_number(s, 6, TOTAL)


# ── SLIDE 7: Services & Menu ───────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
gold_bar(s)
section_label(s, "CMS — Services")

add_text(s, "🍽️  SERVICES, MENU & PRICING", Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         font_size=30, bold=True, color=WHITE)
add_rect(s, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.05), GOLD)

items = [
    ("Event Pricing 💡 NEW", "THE PRICE ENGINE",
     "Set the base per-person rates for Corporate, Tournament, and Birthday events.\nThe Budget Estimator on the Event Place page reads these values in real time.\nUpdate once → all visitor calculators re-calibrate automatically."),
    ("Food Menu", "DIGITAL MENU BOARD",
     "Manage categories (Espresso Bar, Snacks, Mocktails) and individual item prices.\nThe live coffee menu page pulls directly from this module."),
    ("Membership Plans", "TIER EDITOR",
     "Control the features listed for Bronze, Silver, and Gold tiers.\nChanges reflect on the metallic pricing cards visible to new visitors."),
    ("Match My Game", "GAME MATCHING",
     "Manage the skill-level cards that help players find the right competition level."),
]

top = Inches(1.5)
for title, label, desc in items:
    add_rect(s, Inches(0.6), top, Inches(12.1), Inches(1.1), DARK_CARD)
    add_rect(s, Inches(0.6), top, Inches(0.07), Inches(1.1), MAROON)
    add_text(s, title, Inches(0.85), top + Inches(0.06), Inches(2.8), Inches(0.38),
             font_size=13, bold=True, color=WHITE)
    add_text(s, label, Inches(0.85), top + Inches(0.48), Inches(2.8), Inches(0.38),
             font_size=8, bold=True, color=GOLD, italic=True)
    add_text(s, desc, Inches(3.9), top + Inches(0.1), Inches(8.6), Inches(0.85),
             font_size=11, color=GRAY_LIGHT)
    top += Inches(1.2)

slide_number(s, 7, TOTAL)


# ── SLIDE 8: Tournament Bracket & Elimination ──────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
gold_bar(s)
section_label(s, "CMS — Tournaments")

add_text(s, "🏆  TOURNAMENT BRACKETS & ELIMINATION", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         font_size=28, bold=True, color=WHITE)
add_rect(s, Inches(0.8), Inches(1.2), Inches(2.2), Inches(0.05), GOLD)

# Left: steps
steps = [
    ("1", "Generate Bracket", "Choose Gen 4, 8, or 16 to create the tournament skeleton."),
    ("2", "Assign Players", "Drag & Drop — or use Auto Assign for a randomized draw."),
    ("3", "Declare Winners", "Hover on a player → click Trophy Icon → confirm in popup."),
    ("4", "Auto-Advance", "Winner is automatically placed in the correct next-round slot."),
    ("5", "Locked History", "Completed matches are frozen — no accidental edits."),
]

top = Inches(1.55)
for num, title, desc in steps:
    add_rect(s, Inches(0.5), top, Inches(0.5), Inches(0.5), GOLD)
    add_text(s, num, Inches(0.5), top + Inches(0.06), Inches(0.5), Inches(0.38),
             font_size=16, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(1.15), top, Inches(3.0), Inches(0.32),
             font_size=13, bold=True, color=WHITE)
    add_text(s, desc, Inches(1.15), top + Inches(0.3), Inches(5.5), Inches(0.35),
             font_size=11, color=GRAY_LIGHT)
    top += Inches(0.92)

# Right: how it works box
add_rect(s, Inches(7.4), Inches(1.45), Inches(5.5), Inches(5.5), DARK_CARD)
add_rect(s, Inches(7.4), Inches(1.45), Inches(5.5), Inches(0.07), GOLD)
add_text(s, "⚙️  HOW IT WORKS (LAYMAN'S)", Inches(7.6), Inches(1.6), Inches(5.1), Inches(0.45),
         font_size=13, bold=True, color=GOLD)
tech = [
    "Parent-Child Links: Every match 'knows' its next round match.",
    "Winner Push: Clicking Trophy writes the winner ID to the next match.",
    "Smart Slot: Even match → Player 1 slot. Odd match → Player 2 slot.",
    "Permanent Lock: status = 'completed' prevents any future changes.",
    "Champion Toast: Final match win triggers a 🏆 Champion notification.",
]
t = Inches(2.15)
for line in tech:
    add_text(s, "▸  " + line, Inches(7.6), t, Inches(5.0), Inches(0.48),
             font_size=11, color=GRAY_LIGHT)
    t += Inches(0.55)

slide_number(s, 8, TOTAL)


# ── SLIDE 9: Rankings & Security ──────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
gold_bar(s)
section_label(s, "CMS — Rankings & Security")

add_text(s, "📊  RANKINGS & MEMBER SECURITY", Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
         font_size=30, bold=True, color=WHITE)
add_rect(s, Inches(0.8), Inches(1.2), Inches(2.0), Inches(0.05), GOLD)

# Left panel
add_rect(s, Inches(0.5), Inches(1.55), Inches(6.2), Inches(5.4), DARK_CARD)
add_rect(s, Inches(0.5), Inches(1.55), Inches(6.2), Inches(0.07), MAROON)
add_text(s, "🏅  Hall of Fame / Rankings", Inches(0.75), Inches(1.7), Inches(5.8), Inches(0.45),
         font_size=16, bold=True, color=WHITE)
r_items = [
    "Separate leaderboards for Billiards, Darts, and Chess.",
    "Admins can manually add any name — even non-members.",
    "Adding a 'Company' name activates Corporate Rivalry filters on the live site.",
    "Linked accounts show full win/loss stats from tournament play.",
    "Hybrid Linking: Records work with or without a registered user.",
]
t = Inches(2.3)
for it in r_items:
    add_text(s, "▸  " + it, Inches(0.75), t, Inches(5.7), Inches(0.45), font_size=12, color=GRAY_LIGHT)
    t += Inches(0.52)

# Right panel
add_rect(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(5.4), DARK_CARD)
add_rect(s, Inches(7.0), Inches(1.55), Inches(6.0), Inches(0.07), GOLD)
add_text(s, "🛡️  Member Roles & Security", Inches(7.25), Inches(1.7), Inches(5.5), Inches(0.45),
         font_size=16, bold=True, color=WHITE)
s_items = [
    "Every user has a Role Flag: Guest or Admin.",
    "Admin role unlocks all 12 CMS modules invisibly to guests.",
    "Security is enforced server-side (Supabase RLS policies).",
    "Even if someone guesses the URL — Supabase blocks the data.",
    "Promote: Member tab → change tier dropdown to 'Admin'.",
    "Promoted users must refresh their browser to see new tools.",
]
t = Inches(2.3)
for it in s_items:
    add_text(s, "▸  " + it, Inches(7.25), t, Inches(5.5), Inches(0.45), font_size=12, color=GRAY_LIGHT)
    t += Inches(0.52)

slide_number(s, 9, TOTAL)


# ── SLIDE 10: Operations & Contact ────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
gold_bar(s)
section_label(s, "CMS — Operations")

add_text(s, "📅  OPERATIONS & CONTACT", Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         font_size=30, bold=True, color=WHITE)
add_rect(s, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.05), GOLD)

ops = [
    ("📆", "Weekly Schedule", [
        "Assign daily events and special offers (e.g., Industry Night, Family Day).",
        "Each entry has: Day · Event Name · Offer Text · Time slot.",
        "Displayed on the site's timetable automatically, ordered by day.",
    ]),
    ("📞", "Contact Info", [
        "Update address, phone numbers, email, WhatsApp link, and Google Maps URL.",
        "These appear in the footer and the dedicated Contact section sitewide.",
        "Changes propagate to every page instantly after saving.",
    ]),
    ("🔗", "Social Links", [
        "Control Facebook, Instagram, TikTok, and WhatsApp button destinations.",
        "Icons are always visible in the navigation and footer.",
        "Just paste the full URL — no technical knowledge needed.",
    ]),
]

top = Inches(1.5)
for icon, title, bullets in ops:
    add_rect(s, Inches(0.55), top, Inches(12.2), Inches(1.5), DARK_CARD)
    add_rect(s, Inches(0.55), top, Inches(0.08), Inches(1.5), GOLD)
    add_text(s, icon + "  " + title, Inches(0.85), top + Inches(0.1), Inches(3.0), Inches(0.45),
             font_size=15, bold=True, color=GOLD_LIGHT)
    t = top + Inches(0.1)
    for b in bullets:
        add_text(s, "·  " + b, Inches(4.3), t, Inches(8.2), Inches(0.45),
                 font_size=11.5, color=GRAY_LIGHT)
        t += Inches(0.44)
    top += Inches(1.68)

slide_number(s, 10, TOTAL)


# ── SLIDE 11: Database & Migrations ───────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
gold_bar(s)
section_label(s, "Maintenance")

add_text(s, "🧼  DATABASE & MIGRATIONS", Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         font_size=30, bold=True, color=WHITE)
add_rect(s, Inches(0.8), Inches(1.2), Inches(1.5), Inches(0.05), GOLD)

add_text(s, "The database is the foundation. Migrations are like 'building permits' — they officially upgrade the structure.",
         Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5), font_size=14, color=GRAY_LIGHT)

files = [
    ("schema.sql", "Core tables: profiles, players, tournaments, matches, cms_content, registrations."),
    ("cms_expansion.sql", "Adds gallery, food menu, contact, social links, schedule, and pricing CMS modules."),
    ("20260310_bracket_and_event_pricing.sql", "Enables bracket elimination logic + event pricing CMS. Latest migration."),
    ("seed.sql", "Sample data for testing. Safe to run on a fresh database."),
    ("storage_setup.sql", "Creates Supabase Storage buckets for image uploads."),
]

add_text(s, "📁  /supabase/migrations", Inches(0.8), Inches(2.05), Inches(5), Inches(0.4),
         font_size=13, bold=True, color=GOLD)

top = Inches(2.5)
for fname, desc in files:
    add_rect(s, Inches(0.7), top, Inches(12.0), Inches(0.75), DARK_CARD)
    add_text(s, fname, Inches(0.95), top + Inches(0.08), Inches(4.5), Inches(0.3),
             font_size=11, bold=True, color=GOLD_LIGHT)
    add_text(s, desc, Inches(5.6), top + Inches(0.08), Inches(7.0), Inches(0.55),
             font_size=11, color=GRAY_LIGHT)
    top += Inches(0.88)

add_text(s, "To apply: Copy the .sql file → Paste into Supabase SQL Editor → Click Run",
         Inches(0.8), H - Inches(0.55), Inches(11), Inches(0.38),
         font_size=12, bold=True, color=GOLD, italic=True)

slide_number(s, 11, TOTAL)


# ── SLIDE 12: Closing / Thank You ─────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
set_bg(s, DARK_BG)
add_rect(s, 0, 0, Inches(0.18), H, GOLD)
add_rect(s, Inches(0.18), 0, Inches(0.06), H, MAROON)

add_rect(s, Inches(1.2), Inches(3.3), Inches(11), Inches(0.03), GOLD)

add_text(s, "EFREN BILLIARDS & EVENTS", Inches(1.2), Inches(0.9), Inches(11), Inches(1.2),
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s, "Platform Built for Champions 🏆", Inches(1.2), Inches(1.95), Inches(9), Inches(0.7),
         font_size=22, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

add_text(s, "Live at: efrenbilliards.com  ·  Admin Portal: efrenbilliards.com/#admin-cms",
         Inches(1.2), Inches(3.5), Inches(11), Inches(0.5),
         font_size=15, color=GRAY_LIGHT, italic=True)

add_text(s, "For technical support — refer to ADMIN_CMS_GUIDE.md, README.md, and PAGES.md in the project root.",
         Inches(1.2), Inches(4.2), Inches(11), Inches(0.8),
         font_size=13, color=GRAY)

add_text(s, "React  ·  TypeScript  ·  Supabase  ·  Tailwind CSS  ·  Framer Motion",
         Inches(1.2), H - Inches(0.7), Inches(10), Inches(0.4),
         font_size=11, color=GRAY, italic=True)

slide_number(s, 12, TOTAL)


# ── Save ───────────────────────────────────────────────────────────────────
OUT = r"C:\Users\Admin\.gemini\antigravity\scratch\efren-billiards-and-events-place\Efren_Billiards_Platform_Overview.pptx"
prs.save(OUT)
print(f"✅  Saved to: {OUT}")
